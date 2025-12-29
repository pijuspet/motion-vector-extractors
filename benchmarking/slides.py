from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def save_to_ppt(slides, ppt_filename, plots_folder):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout

        # Remove default title placeholder to avoid overlap
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                slide.shapes._spTree.remove(shape._element)

        # Add custom full-width title text box, left aligned
        title_box = slide.shapes.add_textbox(
            0, Inches(0.15), prs.slide_width, Inches(1)
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.text = slide_data["title"]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT

        if slide_data.get("subtitle"):
            left = Inches(0.5)
            top = Inches(1.1)
            width = prs.slide_width - Inches(1)
            height = Inches(0.6)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_box.text_frame
            tf.text = slide_data["subtitle"]
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        img_path = os.path.join(plots_folder, slide_data["filename"])
        if os.path.isfile(img_path):
            left_img = Inches(0.5)
            top_img = Inches(1.9)
            max_width = prs.slide_width - Inches(1)
            max_height = prs.slide_height - Inches(2.1)
            slide.shapes.add_picture(
                img_path, left_img, top_img, width=max_width, height=max_height
            )
    ppt_out = os.path.join(plots_folder, ppt_filename)
    prs.save(ppt_out)
    print(f"\nPowerPoint file created: {ppt_out}")